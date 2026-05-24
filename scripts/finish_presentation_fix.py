# -*- coding: utf-8 -*-
"""Fresh rebuild of the final panorama presentation.

Rules for this pass:
- Use report assets and already generated panorama-script outputs directly.
- Do not fabricate non-equation images or contact sheets.
- Keep missing non-equation visuals as a log entry, not a fake picture.
"""

from __future__ import annotations

import csv
import json
import re
import shutil
import sys
from collections import Counter
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FINAL = ROOT / "final_report"
PRES_DIR = FINAL / "presentation"
PPTX = PRES_DIR / "panorama_stitching_presentation_repaired.pptx"
FALLBACK_PPTX = PRES_DIR / "panorama_stitching_presentation_equation_fixed.pptx"
BACKUP = PRES_DIR / "panorama_stitching_presentation_before_fresh_scan.pptx"
MISSING_LOG = PRES_DIR / "fresh_scan_missing_assets.txt"
EQUATION_DIR = PRES_DIR / "equation_assets"
REPORT_IMG = FINAL / "Computer_Vision___Final_Report" / "asset" / "images"
OUTPUTS = ROOT / "outputs"
DATA_SPLIT = ROOT / "data" / "split"

SLIDE_W, SLIDE_H = Inches(13.333333), Inches(7.5)
BG = RGBColor(248, 250, 252)
INK = RGBColor(20, 25, 35)
MUTED = RGBColor(83, 94, 112)
FAINT = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(0, 122, 128)
BLUE = RGBColor(50, 96, 170)
CORAL = RGBColor(210, 87, 65)
GOLD = RGBColor(194, 139, 35)
GREEN = RGBColor(61, 134, 83)
SLATE = RGBColor(35, 44, 58)

MISSING: list[str] = []


def I(x: float) -> int:
    return Inches(x)


def exists(path: Path | str, label: str = "") -> Path | None:
    p = Path(path)
    if p.exists():
        return p
    item = f"{label}: {p}" if label else str(p)
    if item not in MISSING:
        MISSING.append(item)
    return None


def rows(path: Path) -> list[dict[str, str]]:
    p = exists(path, "csv")
    if p is None:
        return []
    with p.open(newline="", encoding="utf-8-sig") as f:
        return list(csv.DictReader(f))


def img_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as im:
        return ImageOps.exif_transpose(im).size


def bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def text(slide, x, y, w, h, s, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Inches(0.04))
    p = tf.paragraphs[0]
    p.text = str(s)
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def para(slide, x, y, w, h, lines, size=12.5, color=MUTED, bullet=True):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {line}" if bullet else line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def card(slide, x, y, w, h, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = FAINT
    shape.line.width = Pt(1)
    return shape


def chip(slide, x, y, s, color=TEAL, w=None):
    width = w if w is not None else max(0.82, min(3.2, 0.11 * len(s) + 0.35))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(width), I(0.30))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.text = s
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return shape


def title(slide, s, sub=None, section=None):
    if section:
        text(slide, 0.58, 0.20, 2.8, 0.25, section.upper(), 8.5, TEAL, True)
    text(slide, 0.58, 0.42, 10.4, 0.52, s, 24, INK, True)
    if sub:
        text(slide, 0.60, 0.91, 11.3, 0.36, sub, 11, MUTED)


def footer(slide, no):
    text(slide, 0.60, 7.08, 4.2, 0.20, "Panorama Stitching Project", 7.5, MUTED)
    text(slide, 12.20, 7.08, 0.55, 0.20, f"{no:02d}", 7.5, MUTED, align=PP_ALIGN.RIGHT)


def missing_box(slide, path, x, y, w, h, label):
    card(slide, x, y, w, h)
    text(slide, x + 0.15, y + 0.18, w - 0.30, 0.30, "Missing generated asset", 12, CORAL, True)
    text(slide, x + 0.15, y + 0.58, w - 0.30, h - 0.72, f"{label}\n{path}", 8.2, MUTED)


def picture(slide, path, x, y, w, h, label="", border=True, rotation=0):
    p = exists(path, label)
    if p is None:
        missing_box(slide, path, x, y, w, h, label)
        return None
    if border:
        card(slide, x, y, w, h)
    iw, ih = img_size(p)
    rot = rotation % 360
    vw, vh = (ih, iw) if rot in (90, 270) else (iw, ih)
    bw, bh = I(w), I(h)
    scale = min(bw / vw, bh / vh)
    vtw, vth = int(vw * scale), int(vh * scale)
    pw, ph = (vth, vtw) if rot in (90, 270) else (vtw, vth)
    cx, cy = I(x) + bw // 2, I(y) + bh // 2
    pic = slide.shapes.add_picture(str(p), cx - pw // 2, cy - ph // 2, pw, ph)
    if rot:
        pic.rotation = rot
    return pic


def labeled_pic(slide, path, x, y, w, h, label, sub="", rotation=0):
    card(slide, x, y, w, h)
    band_h = 0.42 if sub else 0.30
    picture(slide, path, x + 0.08, y + 0.08, w - 0.16, h - band_h - 0.16, label, False, rotation)
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x + 0.08), I(y + h - band_h - 0.04), I(w - 0.16), I(band_h))
    band.fill.solid()
    band.fill.fore_color.rgb = SLATE
    band.line.fill.background()
    text(slide, x + 0.17, y + h - band_h + 0.03, w - 0.34, 0.15, label, 7.4, WHITE, True)
    if sub:
        text(slide, x + 0.17, y + h - 0.20, w - 0.34, 0.14, sub, 6.1, FAINT)


def grid(slide, items, x, y, w, h, cols, gap=0.12):
    rows_n = (len(items) + cols - 1) // cols
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows_n - 1)) / rows_n
    for idx, item in enumerate(items):
        cx = x + (idx % cols) * (cw + gap)
        cy = y + (idx // cols) * (ch + gap)
        labeled_pic(slide, item["path"], cx, cy, cw, ch, item.get("label", ""), item.get("sub", ""), item.get("rotation", 0))


def stat(slide, x, y, value, label, color=TEAL):
    card(slide, x, y, 1.72, 0.82)
    text(slide, x + 0.08, y + 0.11, 1.55, 0.30, value, 20, color, True, PP_ALIGN.CENTER)
    text(slide, x + 0.08, y + 0.48, 1.55, 0.22, label, 7.4, MUTED, align=PP_ALIGN.CENTER)


def slug(text_value: str) -> str:
    value = re.sub(r"[^a-zA-Z0-9]+", "_", text_value.lower()).strip("_")
    return value[:72] or "equation"


def render_equation_image(name: str, equation_lines) -> Path:
    """Render mathtext equations to a PNG asset.

    These are the only generated presentation images in this rebuild; they are
    mathematical equation assets requested for the theory slides.
    """
    EQUATION_DIR.mkdir(parents=True, exist_ok=True)
    out = EQUATION_DIR / f"{slug(name)}.png"
    lines = equation_lines if isinstance(equation_lines, (list, tuple)) else [equation_lines]
    height = max(1.0, 0.58 * len(lines) + 0.36)
    fig, ax = plt.subplots(figsize=(8.4, height), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_axis_off()
    if len(lines) == 1:
        y_positions = [0.50]
    else:
        top, bottom = 0.78, 0.24
        step = (top - bottom) / max(1, len(lines) - 1)
        y_positions = [top - idx * step for idx in range(len(lines))]
    for line, y_pos in zip(lines, y_positions):
        ax.text(
            0.5,
            y_pos,
            f"${line}$",
            ha="center",
            va="center",
            fontsize=25,
            color="#141923",
            transform=ax.transAxes,
        )
    fig.savefig(out, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    return out


def eq_card(slide, x, y, w, h, name, equation, meaning, symbols, color=TEAL):
    card(slide, x, y, w, h)
    chip(slide, x + 0.16, y + 0.16, "EQUATION", color, 0.95)
    text(slide, x + 1.22, y + 0.13, w - 1.42, 0.26, name, 12.6, INK, True)
    equation_image = render_equation_image(name, equation)
    picture(slide, equation_image, x + 0.30, y + 0.53, w - 0.60, 0.76, name, False)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x + 0.18), I(y + 1.36), I(w - 0.36), I(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = FAINT
    line.line.fill.background()
    text(slide, x + 0.18, y + 1.47, w - 0.36, 0.38, meaning, 8.4, MUTED)
    text(slide, x + 0.18, y + h - 0.43, w - 0.36, 0.22, symbols, 7.0, MUTED)


def table(slide, x, y, w, h, headers, data, color=TEAL, font_size=7.0):
    tbl = slide.shapes.add_table(len(data) + 1, len(headers), I(x), I(y), I(w), I(h)).table
    for ci, head in enumerate(headers):
        c = tbl.cell(0, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.text = head
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for ri, row in enumerate(data, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(245, 247, 250)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                p.font.name = "Arial"
                p.font.size = Pt(font_size)
                p.font.color.rgb = INK
    return tbl


def split_counts():
    out = {}
    for split in ["development", "test", "failure_analysis"]:
        root = DATA_SPLIT / split
        scenes = [p for p in root.iterdir() if p.is_dir() and p.name.startswith("scene_") and (p / "meta.json").exists()] if root.exists() else []
        imgs = sum(len([p for p in s.iterdir() if p.suffix.lower() in {".jpg", ".jpeg", ".png"}]) for s in scenes)
        out[split] = (len(scenes), imgs)
    return out


def scene_source_image(split: str, scene_id: str, index: int = 0) -> Path:
    scene = scene_id if scene_id.startswith("scene_") else f"scene_{scene_id}"
    scene_dir = DATA_SPLIT / split / scene
    meta_path = scene_dir / "meta.json"
    if meta_path.exists():
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8"))
            ordered = meta.get("ordered_files", [])
            if ordered:
                return scene_dir / ordered[min(index, len(ordered) - 1)]
        except json.JSONDecodeError:
            pass
    images = sorted([p for p in scene_dir.glob("*") if p.suffix.lower() in {".jpg", ".jpeg", ".png"}], key=lambda p: p.name.lower())
    return images[min(index, len(images) - 1)] if images else scene_dir / "img_01.jpg"


def method_rows(split):
    out = []
    for r in rows(OUTPUTS / "batch_feature_matching" / "method_summary.csv"):
        if r.get("split") != split:
            continue
        out.append([
            r["method"].replace("HARRIS_HOG", "Harris+HOG"),
            r["pairs"],
            f"{float(r['median_good_matches']):.0f}",
            f"{float(r['median_inliers']):.0f}",
            f"{100 * float(r['median_inlier_ratio']):.1f}%",
            f"{float(r['median_reprojection_error']):.2f}px",
            f"{float(r['median_runtime_sec']):.2f}s",
            f"{100 * float(r['usable_rate']):.1f}%",
        ])
    return out


def comparison_rows():
    out = []
    for r in rows(OUTPUTS / "manual_homography_stitcher" / "comparison" / "manual_vs_opencv_comparison.csv"):
        out.append([
            r["split"].replace("failure_analysis", "failure"),
            r["scene_id"],
            r["comparison_status"],
            f"{r['manual_image_count']}/{r['opencv_num_images']}",
            r["opencv_status_name"],
            f"{100 * float(r['manual_mean_inlier_ratio']):.1f}%" if r.get("manual_mean_inlier_ratio") else "",
            f"{float(r['manual_mean_reprojection_error']):.2f}" if r.get("manual_mean_reprojection_error") else "",
        ])
    return out


def section(prs, no, roman, heading, sub):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, RGBColor(236, 244, 246))
    text(slide, 0.75, 0.78, 1.4, 0.44, roman, 18, TEAL, True)
    text(slide, 0.75, 1.32, 8.8, 0.74, heading, 32, INK, True)
    text(slide, 0.78, 2.18, 9.7, 0.50, sub, 14, MUTED)
    footer(slide, no)


def img_slide(prs, no, heading, sub, path, sec="results", chips=()):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, heading, sub, sec)
    picture(slide, path, 0.68, 1.32, 12.08, 5.12, heading)
    x = 0.86
    for label, color, width in chips:
        chip(slide, x, 6.58, label, color, width)
        x += width + 0.16
    footer(slide, no)


def two_img_slide(prs, no, heading, sub, left, right, left_label, right_label, sec="results"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, heading, sub, sec)
    picture(slide, left, 0.70, 1.34, 5.80, 4.72, left_label)
    picture(slide, right, 6.78, 1.34, 5.72, 4.72, right_label)
    chip(slide, 0.88, 6.28, left_label, TEAL, min(2.5, max(1.1, len(left_label) * 0.12)))
    chip(slide, 7.00, 6.28, right_label, BLUE, min(2.5, max(1.1, len(right_label) * 0.12)))
    footer(slide, no)


def grid_slide(prs, no, heading, sub, items, cols, sec="results"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, heading, sub, sec)
    grid(slide, items, 0.64, 1.28, 12.08, 5.46, cols)
    footer(slide, no)


EQUATION_SLIDES = [
    ("Pipeline toán học tổng quát", "Ảnh -> tiền xử lý -> đặc trưng -> matching/RANSAC -> transform -> panorama", [
        ("Feature-based pipeline", [r"I \rightarrow \tilde{I} \rightarrow (x_i,d_i) \rightarrow M \rightarrow H \rightarrow P"], "Luồng chính từ ảnh đầu vào đến panorama P.", "M: matches; H: transform; P: panorama.", TEAL),
        ("Homography model", [r"x' \sim Hx,\quad x=[u,v,1]^T"], "Ánh xạ hai ảnh khi cảnh gần phẳng hoặc camera quay quanh tâm quang học.", "Dấu ~ nghĩa là bằng nhau tới scale.", BLUE),
        ("Projection to pixels", [r"u'=\frac{x'_1}{x'_3}", r"v'=\frac{x'_2}{x'_3}"], "Chuyển tọa độ homogeneous về pixel để đo lỗi chiếu lại.", "Dùng trong reprojection error.", GOLD),
        ("Transform limits", [r"H \approx \mathrm{single\ global\ model}"], "Parallax, moving objects, weak overlap and sideways scans phá vỡ một transform toàn cục.", "Lý do cần failure_analysis.", CORAL),
    ]),
    ("Preprocessing equations", "Chuẩn hóa kích thước, nhiễu và tương phản trước feature extraction", [
        ("Aspect-preserving resize", [r"s=\min\left(1,\frac{L_{\max}}{\max(W,H)}\right)", r"W'=sW,\quad H'=sH"], "Giảm runtime và bộ nhớ nhưng giữ tỉ lệ hình học.", "L_max là giới hạn cạnh dài.", TEAL),
        ("Gaussian kernel", [r"G_\sigma(x,y)=\frac{1}{2\pi\sigma^2}\exp\left(-\frac{x^2+y^2}{2\sigma^2}\right)"], "Làm mượt nhẹ để giảm nhiễu cao tần.", "sigma điều khiển mức blur.", BLUE),
        ("Gaussian smoothing", [r"\tilde{I}=G_\sigma * I"], "Tạo ảnh feature-ready cho detector/descriptor.", "* là convolution.", GOLD),
        ("CLAHE clip limit", [r"h_c(k)=\min(h(k),T)"], "Tăng tương phản cục bộ nhưng hạn chế khuếch đại nhiễu.", "h(k): histogram tile; T: clip limit.", CORAL),
    ]),
    ("Gradient và Harris Corner", "Cạnh, góc và texture là tín hiệu đầu tiên của matching", [
        ("Gradient magnitude/direction", [r"m=\sqrt{I_x^2+I_y^2}", r"\theta=\mathrm{atan2}(I_y,I_x)"], "Đổi đạo hàm ảnh thành độ mạnh cạnh và hướng.", "m: strength; theta: orientation.", TEAL),
        ("Harris tensor", [r"M=\sum_w [I_x^2,\ I_xI_y;\ I_xI_y,\ I_y^2]"], "Tóm tắt biến thiên gradient trong cửa sổ nhỏ.", "Góc biến thiên mạnh theo hai hướng.", BLUE),
        ("Harris response", [r"R=\det(M)-k\,\mathrm{trace}(M)^2"], "Chấm điểm corner: R lớn dương là góc ổn định.", "Harris dùng làm detector.", GOLD),
        ("HOG bin", [r"h_b=\sum m(x,y)\,\mathbf{1}[\theta(x,y)\in b]"], "Harris+HOG mô tả vùng quanh corner bằng histogram hướng gradient.", "Mỗi bin cộng magnitude.", CORAL),
    ]),
    ("SIFT equations", "Scale-space, DoG và descriptor gradient ổn định nhưng chậm hơn", [
        ("Scale-space", [r"L(x,y,\sigma)=G(x,y,\sigma)*I(x,y)"], "Tạo nhiều mức blur để keypoint sống qua thay đổi scale.", "L là ảnh ở scale sigma.", TEAL),
        ("Difference of Gaussian", [r"D(x,y,\sigma)=L(x,y,k\sigma)-L(x,y,\sigma)"], "Tìm cực trị scale-space làm ứng viên keypoint.", "k là bước scale.", BLUE),
        ("Descriptor normalization", [r"\hat{d}=\frac{d}{\|d\|_2}"], "Giảm nhạy với thay đổi sáng tuyến tính.", "d là vector descriptor.", GOLD),
        ("Report role", [r"e_{\mathrm{reproj}}\downarrow,\quad t_{\mathrm{runtime}}\uparrow"], "SIFT thường cho sai số chiếu lại thấp nhất, đổi lại runtime cao.", "Kết luận từ method_summary.", CORAL),
    ]),
    ("ORB và AKAZE equations", "Binary descriptors nhanh; AKAZE giữ biên tốt trong scale-space phi tuyến", [
        ("ORB orientation", [r"m_{pq}=\sum_x\sum_y x^p y^q I(x,y)", r"\theta=\mathrm{atan2}(m_{01},m_{10})"], "Gán hướng bằng intensity centroid trước khi xoay BRIEF.", "m_pq là moment quanh keypoint.", TEAL),
        ("BRIEF binary test", [r"b_i=\mathbf{1}[I(p_i)<I(q_i)]"], "Tạo từng bit descriptor bằng so sánh hai mẫu cường độ.", "p_i, q_i là sample locations.", BLUE),
        ("Hamming distance", [r"d_H(a,b)=\mathrm{popcount}(a\ \mathrm{XOR}\ b)"], "Đo khác biệt giữa descriptor nhị phân ORB/AKAZE.", "Càng nhỏ càng giống nhau.", GOLD),
        ("AKAZE diffusion", [r"\frac{\partial L}{\partial t}=\mathrm{div}\left(c(x,y,t)\nabla L\right)"], "Scale-space phi tuyến giúp giữ biên tốt hơn Gaussian.", "c điều khiển khuếch tán.", CORAL),
    ]),
    ("Matching, Lowe ratio và RANSAC", "Descriptor-level filtering trước khi kiểm tra hình học", [
        ("Euclidean distance", [r"d_2(a,b)=\sqrt{\sum_j(a_j-b_j)^2}"], "Dùng cho descriptor số thực như SIFT và Harris+HOG.", "Nearest neighbor lấy khoảng cách nhỏ nhất.", TEAL),
        ("Lowe ratio test", [r"d_{\mathrm{best}}<\tau d_{\mathrm{second}},\quad \tau=0.8"], "Loại match mơ hồ, đặc biệt ở repeated patterns.", "best phải rõ ràng hơn second-best.", BLUE),
        ("RANSAC inlier test", [r"\left\|x'_i-\pi(Hx_i)\right\|_2<\theta"], "Giữ match đồng thuận với transform hình học.", "theta là ngưỡng pixel.", GOLD),
        ("RANSAC objective", [r"H^*=\arg\max_H\sum_i \mathbf{1}[e_i(H)<\theta]"], "Chọn H có nhiều inlier nhất qua nhiều mẫu ngẫu nhiên.", "Lỗi thấp chưa đủ nếu inlier hẹp.", CORAL),
    ]),
    ("Warping, canvas, blending và evaluation", "Từ transform hình học sang panorama cuối cùng", [
        ("Inverse warping", [r"x_{\mathrm{src}}\sim H^{-1}x_{\mathrm{dst}}"], "Với mỗi pixel canvas, tìm vị trí nguồn để tránh lỗ warp.", "Dùng trong manual stitcher.", TEAL),
        ("Canvas bbox", [r"\mathrm{bbox}=\min/\max\left(H_k\,\mathrm{corners}(I_k)\right)"], "Tính canvas đủ chứa các góc ảnh sau transform chain.", "Canvas lớn là dấu hiệu drift.", BLUE),
        ("Blending", [r"P(x)=\frac{\sum_k w_k(x)I_k(x)}{\sum_k w_k(x)}"], "Trộn overlap để giảm seam.", "OpenCV có seam/exposure tối ưu hơn.", GOLD),
        ("Coverage/NCC", [r"\mathrm{coverage}=\frac{\mathrm{area}(\mathrm{hull})}{\mathrm{area}(\mathrm{image})}", r"\mathrm{NCC}\ \mathrm{in\ overlap}"], "Đánh giá inlier phân bố rộng và overlap tương đồng.", "Bổ sung cho inlier ratio.", CORAL),
    ]),
]


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    no = 1
    counts = split_counts()
    total_scenes = sum(v[0] for v in counts.values()) or 31
    total_images = sum(v[1] for v in counts.values()) or 167
    showcase_scenes = counts.get("test", (0, 0))[0] + counts.get("failure_analysis", (0, 0))[0]
    comp = rows(OUTPUTS / "manual_homography_stitcher" / "comparison" / "manual_vs_opencv_comparison.csv")
    comp_counts = Counter(r["comparison_status"] for r in comp)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    picture(slide, OUTPUTS / "openCV" / "panoramas" / "scene_01_opencv_panorama.jpg", 0.58, 0.46, 12.18, 3.38, "cover panorama")
    text(slide, 0.78, 4.12, 8.8, 0.34, "COMPUTER VISION - BÁO CÁO THI CUỐI KỲ", 9.2, TEAL, True)
    text(slide, 0.75, 4.54, 9.8, 0.80, "Xây dựng ảnh toàn cảnh", 36, INK, True)
    text(slide, 0.78, 5.38, 10.6, 0.36, "Panorama stitching pipeline: dữ liệu, lý thuyết, matching/RANSAC, OpenCV baseline và manual affine-chain.", 13, MUTED)
    text(slide, 0.78, 6.18, 8.8, 0.50, "Phạm Hùng Sơn - Nguyễn Minh Quân - Trần Huy Giang\nNhóm 15 - HUS / VNU", 11, SLATE)
    footer(slide, no)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Flow thuyết trình theo báo cáo", "Mở đầu -> Cơ sở lý thuyết -> Dữ liệu -> Kết quả -> Android -> Kết luận", "overview")
    for idx, (num, head, body, col) in enumerate([
        ("01", "Vấn đề", "Vì sao panorama stitching cần căn chỉnh hình học và kiểm soát dữ liệu.", TEAL),
        ("02", "Lý thuyết", "Preprocessing, feature, matching, RANSAC, homography, warping, blending.", BLUE),
        ("03", "Thực nghiệm", "Scene examples, failure diversity, batch matching, OpenCV và manual stitcher.", GOLD),
        ("04", "Ứng dụng", "Portable pipeline và Android offline bridge theo nội dung báo cáo.", CORAL),
    ]):
        x = 0.74 + idx * 3.05
        card(slide, x, 1.72, 2.58, 3.86)
        chip(slide, x + 0.18, 1.95, num, col, 0.56)
        text(slide, x + 0.18, 2.45, 2.18, 0.34, head, 17, INK, True)
        text(slide, x + 0.18, 3.04, 2.18, 1.42, body, 11.3, MUTED)
    footer(slide, no)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "I. Giới thiệu vấn đề", "Từ nhiều ảnh chồng lấp đến một panorama góc rộng ổn định", "problem")
    para(slide, 0.78, 1.45, 5.55, 3.90, [
        "Đầu vào là chuỗi ảnh cùng cảnh, có overlap nhưng khác viewpoint, ánh sáng và độ sắc nét.",
        "Mục tiêu là tìm hệ tọa độ chung, warp ảnh vào canvas và blend vùng chồng lấp.",
        "Báo cáo phân tích success, hard-valid và failure để giải thích vì sao pipeline hoạt động hoặc đứt chuỗi.",
        "Lỗi chính: blur, ít texture, mẫu lặp, thiếu overlap, parallax, vật thể chuyển động và exposure change.",
    ], 14)
    picture(slide, REPORT_IMG / "pipeline.png", 6.74, 1.38, 5.65, 4.10, "report pipeline")
    footer(slide, no)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Project goals", "Build panoramas and explain why each scene succeeds, becomes difficult, or fails", "problem")
    for x, value, label, col in [(0.76, total_scenes, "curated scenes", TEAL), (2.70, total_images, "input images", BLUE), (4.64, 4, "main descriptors", GOLD), (6.58, showcase_scenes, "showcase scenes", CORAL)]:
        stat(slide, x, 1.42, str(value), label, col)
    para(slide, 0.82, 2.65, 5.78, 2.96, [
        "Tổ chức dữ liệu theo data/split và kiểm định đầu vào bằng metadata/audit.",
        "So sánh ORB, AKAZE, Harris+HOG và SIFT bằng match, inlier, reprojection error, coverage và runtime.",
        "Xây dựng manual geometry stitcher; báo cáo cuối dùng affine mặc định để giảm over-warp khi chain nhiều ảnh điện thoại.",
    ], 12.8)
    para(slide, 7.06, 2.65, 5.20, 2.96, [
        "So sánh manual stitcher với cv2.Stitcher ở mức scene.",
        "Phân biệt both_ok, manual_only và manual_partial_only.",
        "Giữ failure artifacts như bằng chứng để phân tích giới hạn, không chỉ trình bày scene đẹp.",
        "Scene IDs are sparse after pruning; scene_40 is an identifier, not a dataset count.",
    ], 12.8)
    footer(slide, no)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "II. Dataset & experiments", f"The final curated dataset is data/split: {total_scenes} scenes, {total_images} images, sparse scene IDs", "data")
    for idx, (split, desc, col) in enumerate([("development", "reserve/tuning", TEAL), ("test", "showcase hợp lệ", BLUE), ("failure_analysis", "failure và hard-valid", CORAL)]):
        x = 0.82 + idx * 4.08
        scenes, imgs = counts[split]
        card(slide, x, 1.44, 3.50, 1.72)
        chip(slide, x + 0.18, 1.65, split, col, 1.56)
        text(slide, x + 0.18, 2.08, 1.35, 0.42, str(scenes), 25, col, True)
        text(slide, x + 1.34, 2.18, 0.72, 0.22, "scenes", 8.2, MUTED)
        text(slide, x + 2.05, 2.08, 0.88, 0.42, str(imgs), 25, col, True)
        text(slide, x + 2.78, 2.18, 0.56, 0.22, "images", 8.2, MUTED)
        text(slide, x + 0.18, 2.68, 3.06, 0.20, desc, 8.8, MUTED)
    picture(slide, OUTPUTS / "batch_feature_matching" / "plots" / "scene_method_usable_rate_heatmap.png", 0.80, 3.55, 5.85, 2.78, "usable-rate heatmap")
    picture(slide, OUTPUTS / "batch_feature_matching" / "plots" / "method_quality_bubble_by_split.png", 6.86, 3.55, 5.62, 2.78, "method quality bubble")
    footer(slide, no)
    no += 1

    grid_slide(prs, no, "Scene example: success", "Ảnh gốc từ report asset, đặt contain-fit để không bị crop", [
        {"path": REPORT_IMG / "scene_example" / "success-example" / f"img_{i:02d}.jpg", "label": f"img_{i:02d}"}
        for i in range(1, 7)
    ], 3, "data")
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Scene examples: hard-valid and failure", "Portrait source photos are rotated in PowerPoint; no new image assets are generated", "data")
    grid(slide, [{"path": REPORT_IMG / "scene_example" / "hard-valid-example" / f"img_{i:02d}.jpg", "label": f"hard-valid {i}", "rotation": 270} for i in range(1, 6)], 0.68, 1.30, 7.56, 4.78, 2)
    grid(slide, [{"path": REPORT_IMG / "scene_example" / "failure-example" / f"img_{i:02d}.jpg", "label": f"failure {i}", "rotation": 270} for i in range(1, 4)], 8.46, 1.30, 3.96, 4.78, 1)
    chip(slide, 0.86, 6.36, "hard-valid", GOLD, 1.35)
    chip(slide, 2.38, 6.36, "failure", CORAL, 1.0)
    footer(slide, no)
    no += 1

    div = [
        ("03", "wrong image order"),
        ("04", "insufficient overlap"),
        ("06", "sky dominance"),
        ("07", "zoom / scale change"),
        ("08", "global stitch failure"),
        ("11", "low texture"),
        ("15", "weak global connectivity"),
        ("16", "translated camera motion"),
        ("18", "low light / exposure boost"),
        ("21", "wide sweep / exposure"),
        ("30", "output variation"),
        ("32", "repeated pattern"),
        ("35", "motion blur / moving objects"),
    ]
    for chunk_idx, chunk in enumerate([div[:7], div[7:]], start=1):
        grid_slide(prs, no, f"Failure scene diversity showcase {chunk_idx}", "Current failure_analysis split; source frames are used directly without generated contact sheets", [
            {"path": scene_source_image("failure_analysis", sid), "label": f"scene_{sid}", "sub": label}
            for sid, label in chunk
        ], 4, "data")
        no += 1

    section(prs, no, "III", "Cơ sở lý thuyết", "Mỗi equation được tách ra cùng ý nghĩa và vai trò trong pipeline.")
    no += 1
    for heading, sub, eqs in EQUATION_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide)
        title(slide, heading, sub, "theory")
        positions = [(0.72, 1.32), (6.78, 1.32), (0.72, 3.98), (6.78, 3.98)]
        for (name, equation, meaning, symbols, col), (x, y) in zip(eqs, positions):
            eq_card(slide, x, y, 5.85, 2.26, name, equation, meaning, symbols, col)
        footer(slide, no)
        no += 1

    section(prs, no, "IV", "Result showcase", "Các hình sau dùng trực tiếp từ report assets hoặc outputs do notebook/script tạo.")
    no += 1
    img_slide(prs, no, "Kết quả sau tiền xử lý", "Report asset: resize -> grayscale -> Gaussian blur -> CLAHE", REPORT_IMG / "preprocessing_result" / "preprocessing-full-pipeline.png", "results", [("existing report figure", TEAL, 1.90), ("no regenerated image", BLUE, 1.92)])
    no += 1

    mv = OUTPUTS / "case_study" / "failure_analysis" / "scene_18" / "pair_01" / "method_visualizations"
    grid_slide(prs, no, "Detector và descriptor visualizations", "Generated by Notebook 03 / case_study outputs for scene_18 pair_01", [
        {"path": mv / "01_detector_cues.png", "label": "detector cues"},
        {"path": mv / "02_pipeline_keypoint_overlays.png", "label": "keypoint overlays"},
        {"path": mv / "03_descriptor_vector_examples.png", "label": "descriptor vectors"},
        {"path": mv / "04_hog_patch_descriptor.png", "label": "HOG patch"},
    ], 2)
    no += 1

    for case_name, heading, sub, col in [
        ("descriptor-example-success", "Success scene: feature -> match -> stitch", "Clean overlap produces dense, consistent RANSAC inliers.", GREEN),
        ("descriptor-example-hard-valid", "Hard-valid scene: valid but difficult", "Repeated patterns and viewpoint changes make the output sensitive.", GOLD),
        ("descriptor-example-failure", "Failure scene: weak/ambiguous geometry", "The failure artifact is included to explain pipeline limits.", CORAL),
    ]:
        base = REPORT_IMG / "feature_method_result" / case_name
        grid_slide(prs, no, heading, sub, [
            {"path": base / "preprocessed-pair.png", "label": "preprocessed pair"},
            {"path": base / "matching.png", "label": "matching / inliers"},
            {"path": base / "stitching.png", "label": "stitching output"},
        ], 3)
        no += 1

    for split, scene, pair, heading, sub in [
        ("failure_analysis", "scene_18", "pair_01", "Pair panoramas by descriptor", "Same adjacent pair, different feature pipelines: scene_18 pair_01"),
    ]:
        p = OUTPUTS / "case_study" / split / scene / pair / "panoramas"
        grid_slide(prs, no, heading, sub, [
            {"path": p / "sift_panorama.jpg", "label": "SIFT"},
            {"path": p / "orb_panorama.jpg", "label": "ORB"},
            {"path": p / "akaze_panorama.jpg", "label": "AKAZE"},
            {"path": p / "harris_hog_panorama.jpg", "label": "Harris+HOG"},
        ], 2)
        no += 1


    two_img_slide(prs, no, "Batch feature matching overview", "Descriptor reliability across all adjacent pairs", OUTPUTS / "batch_feature_matching" / "plots" / "scene_method_usable_rate_heatmap.png", OUTPUTS / "batch_feature_matching" / "plots" / "method_quality_bubble_by_split.png", "usable-rate heatmap", "quality bubble")
    no += 1
    two_img_slide(prs, no, "Batch metric distributions", "Không chỉ bảng số liệu: phân phối và tương quan inlier/error", OUTPUTS / "batch_feature_matching" / "plots" / "metric_distributions_by_method.png", OUTPUTS / "batch_feature_matching" / "plots" / "pair_inliers_vs_reprojection_error.png", "metric distributions", "inliers vs reprojection")
    no += 1
    grid_slide(prs, no, "Test vs failure_analysis results", "Status and usable-rate visualizations from Notebook 04", [
        {"path": OUTPUTS / "batch_feature_matching" / "plots" / "test_status_distribution_by_method.png", "label": "test status"},
        {"path": OUTPUTS / "batch_feature_matching" / "plots" / "test_usable_rate_by_method.png", "label": "test usable rate"},
        {"path": OUTPUTS / "batch_feature_matching" / "plots" / "failure_analysis_status_distribution_by_method.png", "label": "failure status"},
        {"path": OUTPUTS / "batch_feature_matching" / "plots" / "failure_analysis_usable_rate_by_method.png", "label": "failure usable rate"},
    ], 2)
    no += 1
    two_img_slide(prs, no, "Weak pair diagnostics", "Failure analysis includes pair-level status matrices and metric timelines", OUTPUTS / "batch_feature_matching" / "plots" / "failure_analysis_pair_status_overview.png", OUTPUTS / "batch_feature_matching" / "plots" / "weakest_adjacent_pairs_status_stack.png", "failure pair overview", "weakest pairs")
    no += 1
    two_img_slide(prs, no, "Scene-level diagnostics", "Example matrix and timeline reveal where a sequence breaks", OUTPUTS / "batch_feature_matching" / "plots" / "test_scene_36_pair_status_matrix.png", OUTPUTS / "batch_feature_matching" / "plots" / "failure_analysis_scene_15_pair_metric_timeline.png", "scene_36 matrix", "scene_15 timeline")
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Method summary table", "Numeric summaries from outputs/batch_feature_matching/method_summary.csv", "results")
    headers = ["Method", "Pairs", "Good", "Inliers", "Inlier %", "Reproj.", "Runtime", "Usable"]
    text(slide, 0.78, 1.22, 3.6, 0.26, "Test split", 13, TEAL, True)
    table(slide, 0.72, 1.58, 11.88, 1.52, headers, method_rows("test"), TEAL)
    text(slide, 0.78, 3.44, 3.8, 0.26, "Failure analysis split", 13, CORAL, True)
    table(slide, 0.72, 3.80, 11.88, 1.52, headers, method_rows("failure_analysis"), CORAL)
    para(slide, 0.78, 5.74, 11.7, 0.66, ["Report takeaway: AKAZE and SIFT provide dense/accurate inliers; ORB is fastest; failure_analysis reduces spatial coverage below the clean-scene regime."], 10.2, MUTED, False)
    footer(slide, no)
    no += 1

    grid_slide(prs, no, "OpenCV Stitcher baseline outputs", "Direct panorama files from outputs/openCV/panoramas", [
        {"path": OUTPUTS / "openCV" / "panoramas" / "scene_01_opencv_panorama.jpg", "label": "scene_01", "sub": "test success"},
        {"path": OUTPUTS / "openCV" / "panoramas" / "scene_21_opencv_panorama.jpg", "label": "scene_21", "sub": "wide hard-valid"},
        {"path": OUTPUTS / "openCV" / "panoramas" / "scene_30_opencv_panorama.jpg", "label": "scene_30", "sub": "exposure variation"},
    ], 3)
    no += 1
    two_img_slide(prs, no, "Scene problems affecting stitching", "Problem tags explain why status codes are not enough", OUTPUTS / "failure_problem_analysis" / "failure_problem_status_counts.png", REPORT_IMG / "final_results" / "failure_problem_status_counts.png", "generated output chart", "report figure")
    no += 1
    grid_slide(prs, no, "Manual geometry stitcher vs OpenCV", "Comparison script outputs over the current test and failure_analysis showcase scenes", [
        {"path": OUTPUTS / "manual_homography_stitcher" / "comparison" / "plots" / "manual_vs_opencv_status_counts.png", "label": "status counts"},
        {"path": OUTPUTS / "manual_homography_stitcher" / "comparison" / "plots" / "manual_vs_opencv_area.png", "label": "panorama area"},
        {"path": OUTPUTS / "manual_homography_stitcher" / "comparison" / "plots" / "manual_runtime_by_scene.png", "label": "manual runtime"},
    ], 3)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Scene comparison table", "Rows read from manual_vs_opencv_comparison.csv", "results")
    for x, key, label, col in [
        (0.78, "both_ok", "both_ok", GREEN),
        (2.70, "both_ok_manual_partial", "OpenCV full/manual partial", GOLD),
        (4.62, "manual_only", "manual_only", BLUE),
        (6.54, "manual_partial_only", "manual partial only", CORAL),
    ]:
        stat(slide, x, 1.22, str(comp_counts.get(key, 0)), label, col)
    table(slide, 0.72, 2.36, 11.88, 3.75, ["Split", "Scene", "Status", "Manual imgs", "OpenCV", "Inlier", "Reproj."], comparison_rows(), TEAL)
    footer(slide, no)
    no += 1

    showcase_groups = [
        ("Test showcase I", [("test", "01", "both_ok clean success"), ("test", "24", "manual_only, OpenCV camera fail"), ("test", "28", "OpenCV full, manual partial")]),
        ("Test showcase II", [("test", "33", "OpenCV full, manual partial"), ("test", "36", "both_ok hard-valid"), ("test", "40", "both_ok long sequence")]),
        ("Failure showcase I", [("failure_analysis", "03", "wrong image order"), ("failure_analysis", "04", "insufficient overlap"), ("failure_analysis", "06", "sky dominance")]),
        ("Failure showcase II", [("failure_analysis", "07", "zoom/scale change"), ("failure_analysis", "08", "global stitch failure"), ("failure_analysis", "11", "low texture")]),
        ("Failure showcase III", [("failure_analysis", "15", "weak global connectivity"), ("failure_analysis", "16", "translated camera motion"), ("failure_analysis", "21", "wide sweep/exposure")]),
        ("Failure showcase IV", [("failure_analysis", "30", "output variation"), ("failure_analysis", "32", "repeated pattern"), ("failure_analysis", "35", "motion blur/moving objects")]),
    ]
    for heading, scenes in showcase_groups:
        items = []
        for split, sid, sub in scenes:
            prefix = "failure_analysis" if split == "failure_analysis" else "test"
            items.append({"path": OUTPUTS / "manual_homography_stitcher" / "comparison" / "plots" / f"manual_vs_opencv_{prefix}_scene_{sid}.png", "label": f"scene_{sid}", "sub": sub})
        grid_slide(prs, no, heading, "Side-by-side comparison plots generated by the comparison script", items, 3)
        no += 1

    grid_slide(prs, no, "Portable / mobile pipeline verification", "Generated outputs from mobile_app_pipeline_check and portable_review", [
        {"path": OUTPUTS / "mobile_app_pipeline_check" / "scene_01" / "manual_ORB" / "panorama.jpg", "label": "manual ORB"},
        {"path": OUTPUTS / "mobile_app_pipeline_check" / "scene_01" / "manual_AKAZE" / "panorama.jpg", "label": "manual AKAZE"},
        {"path": OUTPUTS / "mobile_app_pipeline_check" / "scene_01" / "manual_SIFT" / "panorama.jpg", "label": "manual SIFT"},
    ], 2, "android")
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Android app integration", "Report content: offline Chaquopy pipeline and Python-to-JS progress bridge", "android")
    para(slide, 0.78, 1.42, 5.65, 4.54, [
        "Python pipeline packaged into the Android app through Chaquopy.",
        "Real-time progress bridge writes/reads a shared progress file so slow SIFT/Harris+HOG steps do not look frozen.",
        "Student Debug Mode compares OpenCV Stitcher and manual methods with ORB, SIFT, AKAZE, Harris+HOG.",
        "The app exposes metrics: good matches, RANSAC inliers, inlier ratio, reprojection error, coverage, NCC and homography sanity.",
    ], 13.2)
    card(slide, 6.92, 1.54, 5.26, 3.70)
    text(slide, 7.24, 1.88, 4.58, 0.36, "Screenshot assets", 17, INK, True)
    text(slide, 7.24, 2.46, 4.58, 1.48, "No generated Android UI screenshots were found in the project scan, so this deck does not fabricate them.", 14, MUTED)
    text(slide, 7.24, 4.22, 4.58, 0.50, "Inserted instead: existing mobile/portable pipeline output artifacts.", 11.2, TEAL, True)
    footer(slide, no)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title(slide, "Kết luận", "Pipeline đã hoàn thiện từ dữ liệu đến panorama và ứng dụng Android offline", "closing")
    para(slide, 0.82, 1.36, 5.70, 4.46, [
        "Đã hệ thống hóa lý thuyết: preprocessing, detector/descriptor, matching, RANSAC, homography, warping và blending.",
        "Đã xây dựng dataset, audit và batch evaluation cho ORB, AKAZE, Harris+HOG, SIFT.",
        "OpenCV baseline mạnh ở scene hợp lệ; failure_analysis cho thấy giới hạn do parallax, thiếu overlap và scan tịnh tiến.",
    ], 13.4)
    para(slide, 6.92, 1.36, 5.30, 4.46, [
        "Manual affine-chain minh bạch và hữu ích để chẩn đoán, nhưng chưa thay thế OpenCV về bundle adjustment, seam optimization và exposure compensation.",
        "Failure artifacts không bị giấu; chúng là bằng chứng để giải thích giới hạn và hướng phát triển tiếp theo.",
        "Hướng tiếp theo: global optimization, projection tốt hơn, local warping và blending nâng cao.",
    ], 13.4)
    footer(slide, no)
    no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    picture(slide, OUTPUTS / "openCV" / "panoramas" / "scene_40_opencv_panorama.jpg", 0.70, 0.74, 11.90, 3.24, "closing panorama")
    text(slide, 0.84, 4.55, 6.2, 0.78, "Thank you", 38, INK, True)
    text(slide, 0.88, 5.40, 6.2, 0.36, "Q&A", 18, TEAL, True)
    text(slide, 0.88, 6.04, 10.8, 0.32, "Panorama stitching = data quality + local features + robust geometry + careful blending.", 13, MUTED)
    footer(slide, no)
    return prs


def write_missing_log():
    if MISSING:
        msg = "Missing non-equation assets from fresh scan:\n" + "\n".join(f"- {x}" for x in MISSING) + "\n"
    else:
        msg = "No missing non-equation assets were detected during the fresh presentation scan.\n"
    MISSING_LOG.write_text(msg, encoding="utf-8")


def main() -> int:
    PRES_DIR.mkdir(parents=True, exist_ok=True)
    if PPTX.exists() and not BACKUP.exists():
        shutil.copy2(PPTX, BACKUP)
        print(f"Backed up current presentation: {BACKUP}")
    prs = build()
    saved_path = PPTX
    try:
        prs.save(PPTX)
    except PermissionError:
        saved_path = FALLBACK_PPTX
        prs.save(saved_path)
    write_missing_log()
    print(f"Saved fresh presentation: {saved_path}")
    print(f"Slides: {len(prs.slides)}")
    if MISSING:
        print(f"Missing assets logged: {MISSING_LOG}")
        for item in MISSING:
            print(f"MISSING: {item}")
    else:
        print(f"Missing assets: none ({MISSING_LOG})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
